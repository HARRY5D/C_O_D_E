"""
PDF to PowerPoint Converter for Research Papers
Extracts content from PDF and creates a professional presentation
"""

import PyPDF2
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def extract_pdf_content(pdf_path):
    """Extract text content from PDF"""
    content = []
    try:
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            print(f"Total pages: {len(pdf_reader.pages)}")
            
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text.strip():
                    content.append(text)
        
        full_text = "\n".join(content)
        return full_text
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return ""

def parse_research_paper(text):
    """Parse research paper into structured sections"""
    sections = {
        'title': '',
        'abstract': '',
        'introduction': '',
        'methodology': '',
        'results': '',
        'conclusions': '',
        'references': '',
        'key_points': []
    }
    
    lines = text.split('\n')
    
    # Extract title (usually first meaningful line)
    for line in lines[:20]:
        if line.strip() and len(line.strip()) > 10:
            if 'robustness' in line.lower() or 'framework' in line.lower():
                sections['title'] = line.strip()
                break
    
    # Extract sections
    current_section = None
    section_content = []
    
    for line in lines:
        lower_line = line.lower()
        
        if any(keyword in lower_line for keyword in ['abstract', 'summary']):
            if section_content and current_section:
                sections[current_section] = '\n'.join(section_content)
            current_section = 'abstract'
            section_content = []
        elif any(keyword in lower_line for keyword in ['introduction', 'background']):
            if section_content and current_section:
                sections[current_section] = '\n'.join(section_content)
            current_section = 'introduction'
            section_content = []
        elif any(keyword in lower_line for keyword in ['methodology', 'approach', 'method']):
            if section_content and current_section:
                sections[current_section] = '\n'.join(section_content)
            current_section = 'methodology'
            section_content = []
        elif any(keyword in lower_line for keyword in ['result', 'finding', 'evaluation']):
            if section_content and current_section:
                sections[current_section] = '\n'.join(section_content)
            current_section = 'results'
            section_content = []
        elif any(keyword in lower_line for keyword in ['conclusion', 'discussion', 'future work']):
            if section_content and current_section:
                sections[current_section] = '\n'.join(section_content)
            current_section = 'conclusions'
            section_content = []
        elif any(keyword in lower_line for keyword in ['reference', 'bibliography']):
            if section_content and current_section:
                sections[current_section] = '\n'.join(section_content)
            current_section = 'references'
            section_content = []
        
        if current_section and line.strip():
            section_content.append(line.strip())
    
    # Extract key points from content
    full_content = ' '.join([sections[k] for k in sections if k != 'key_points'])
    sentences = re.split(r'[.!?]+', full_content)
    
    for sentence in sentences:
        sentence = sentence.strip()
        if len(sentence) > 20 and any(keyword in sentence.lower() for keyword in 
                                       ['framework', 'robustness', 'testing', 'method', 'approach', 'solution']):
            sections['key_points'].append(sentence)
            if len(sections['key_points']) >= 8:
                break
    
    return sections

def create_presentation(sections, output_path):
    """Create PowerPoint presentation from sections"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    TITLE_COLOR = RGBColor(0, 51, 102)  # Dark blue
    ACCENT_COLOR = RGBColor(0, 102, 204)  # Medium blue
    TEXT_COLOR = RGBColor(64, 64, 64)  # Dark gray
    
    def add_title_slide(title, subtitle):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle if subtitle else "Research Paper Presentation"
        p.font.size = Pt(28)
        p.font.color.rgb = ACCENT_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(title, content_points):
        """Add content slide with bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        # Add line under title
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(0))
        line.line.color.rgb = ACCENT_COLOR
        line.line.width = Pt(3)
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(content_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point.strip()
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(6)
            p.space_after = Pt(6)
    
    # Slide 1: Title Slide
    add_title_slide(
        sections.get('title', 'Robustness Testing Framework'),
        "Research Paper Presentation"
    )
    
    # Slide 2: Abstract
    if sections['abstract']:
        abstract_text = sections['abstract'][:500]
        add_content_slide("Abstract", [abstract_text])
    
    # Slide 3: Introduction
    if sections['introduction']:
        intro_points = [p.strip() for p in sections['introduction'].split('\n') if p.strip()][:5]
        add_content_slide("Introduction", intro_points if intro_points else ["Understanding the need for robust testing frameworks"])
    
    # Slide 4: Key Research Points
    if sections['key_points']:
        add_content_slide("Key Research Points", sections['key_points'][:5])
    
    # Slide 5: Methodology
    if sections['methodology']:
        method_points = [p.strip() for p in sections['methodology'].split('\n') if p.strip()][:5]
        add_content_slide("Methodology", method_points if method_points else ["Research methodology and approach"])
    
    # Slide 6: Results & Findings
    if sections['results']:
        result_points = [p.strip() for p in sections['results'].split('\n') if p.strip()][:5]
        add_content_slide("Results & Findings", result_points if result_points else ["Key findings and results"])
    
    # Slide 7: Conclusions
    if sections['conclusions']:
        conclusion_points = [p.strip() for p in sections['conclusions'].split('\n') if p.strip()][:5]
        add_content_slide("Conclusions", conclusion_points if conclusion_points else ["Key takeaways and conclusions"])
    
    # Slide 8: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = TITLE_COLOR
    
    thank_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    thank_frame = thank_box.text_frame
    p = thank_frame.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save(output_path)
    print(f"Presentation saved: {output_path}")

def main():
    pdf_path = r"D:\COLLEGE\SEMS\6th Sem\Research\Robustness_Testing_Framework.pdf"
    output_dir = r"d:\JAVA\CODE\Crypto_code\mfa-ppt"
    output_path = os.path.join(output_dir, "Robustness_Testing_Framework_Presentation.pptx")
    
    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)
    
    print("Extracting PDF content...")
    pdf_content = extract_pdf_content(pdf_path)
    
    if not pdf_content:
        print("Error: Could not extract content from PDF")
        return
    
    print("Parsing research paper...")
    sections = parse_research_paper(pdf_content)
    
    print("Creating presentation...")
    create_presentation(sections, output_path)
    print("\nPresentation created successfully!")

if __name__ == "__main__":
    main()
